#!/usr/bin/env python3
"""
Create PowerPoint presentation for Network Benchmark Soutenance
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ENSEEIHT colors
ENSEEIHT_BLUE = RGBColor(0, 76, 153)
DARK_GRAY = RGBColor(64, 64, 64)
LIGHT_GRAY = RGBColor(128, 128, 128)
GREEN = RGBColor(0, 153, 76)
RED = RGBColor(204, 0, 0)
ORANGE = RGBColor(255, 153, 0)

def create_title_slide(prs, title, subtitle):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ENSEEIHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title, layout_index=1):
    """Create content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = ENSEEIHT_BLUE
    title_shape.text_frame.paragraphs[0].font.bold = True
    return slide

def add_bullet_points(slide, points, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4)):
    """Add bullet points to slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(10)
    
    return textbox

def add_table(slide, data, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4)):
    """Add table to slide"""
    rows = len(data)
    cols = len(data[0])
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    for col_idx, cell_text in enumerate(data[0]):
        cell = table.rows[0].cells[col_idx]
        cell.text = str(cell_text)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ENSEEIHT_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
    
    # Data rows
    for row_idx in range(1, rows):
        for col_idx, cell_text in enumerate(data[row_idx]):
            cell = table.rows[row_idx].cells[col_idx]
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
    
    return table

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ============================================================================
    # SLIDE 1: Title
    # ============================================================================
    slide = create_title_slide(
        prs,
        "Energy-Efficient Serverless Computing",
        "A Comparative Analysis of VM and MicroVM Isolation"
    )
    
    # Add author info
    author_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    author_frame = author_box.text_frame
    p = author_frame.paragraphs[0]
    p.text = "Ayoub SAMI\nENSEEIHT - Projet Long 2025-2026\nSupervisor: Prof. Daniel Hagimont"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    
    # ============================================================================
    # SLIDE 2: Context & Motivation
    # ============================================================================
    slide = create_content_slide(prs, "Context & Motivation")
    points = [
        "Serverless computing is transforming cloud infrastructure",
        "Security vs Performance trade-off in multi-tenant environments",
        "Traditional containers: Fast but kernel-level isolation only",
        "MicroVMs: Hardware isolation but potentially higher overhead",
        "Research Question: What is the real cost of MicroVM security?"
    ]
    add_bullet_points(slide, points)
    
    # ============================================================================
    # SLIDE 3: Research Questions
    # ============================================================================
    slide = create_content_slide(prs, "Research Questions")
    points = [
        "1. Performance Impact: MicroVM overhead vs containers?",
        "2. Energy Efficiency: How do MicroVMs compare in power consumption?",
        "3. Cold Start Penalty: Impact across isolation mechanisms?",
        "4. Network Performance: Does virtualization degrade I/O?",
        "5. Language Impact: How do runtimes interact with isolation?"
    ]
    add_bullet_points(slide, points)
    
    # ============================================================================
    # SLIDE 4: Experimental Setup
    # ============================================================================
    slide = create_content_slide(prs, "Experimental Infrastructure")
    
    table_data = [
        ["Component", "Cluster 1 (Baseline)", "Cluster 2 (MicroVM)"],
        ["Runtime", "Docker + runc", "Kata Containers + QEMU"],
        ["Kubernetes", "v1.28.2", "v1.28.2"],
        ["Hardware", "Intel Xeon E5-2680 v4", "Intel Xeon E5-2680 v4"],
        ["Memory", "64GB DDR4", "64GB DDR4"],
        ["Network", "10 Gbps Ethernet", "10 Gbps Ethernet"],
        ["FaaS Platform", "OpenWhisk 1.0", "OpenWhisk 1.0"]
    ]
    add_table(slide, table_data, Inches(0.5), Inches(2), Inches(9), Inches(4))
    
    # ============================================================================
    # SLIDE 5: Test Matrix
    # ============================================================================
    slide = create_content_slide(prs, "Benchmark Design")
    points = [
        "Network-intensive workloads across 3 languages:",
        "  • Python 3.9",
        "  • JavaScript (Node.js 14)",
        "  • Java 8",
        "4 network operations per test:",
        "  • TCP Latency (8.8.8.8:53)",
        "  • HTTP Latency (httpbin.org)",
        "  • DNS Resolution (google.com)",
        "  • Bandwidth (42.66 MB download)",
        "Test matrix: 2 clusters × 3 languages × 2 start types × 5 iterations = 60 tests"
    ]
    add_bullet_points(slide, points)
    
    # ============================================================================
    # SLIDE 6: Performance Results - Summary
    # ============================================================================
    slide = create_content_slide(prs, "Performance Results: Summary")
    
    table_data = [
        ["Test", "C1 (VM)", "C2 (µVM)", "Overhead"],
        ["Python Cold", "41.84s", "47.17s", "+12.7%"],
        ["Python Hot", "26.37s", "25.65s", "-2.7%"],
        ["JavaScript Cold", "41.66s", "45.43s", "+9.0%"],
        ["JavaScript Hot", "24.72s", "25.59s", "+3.5%"],
        ["Java Cold", "45.17s", "44.73s", "-1.0%"],
        ["Java Hot", "23.57s", "23.85s", "+1.2%"],
        ["Average", "38.89s", "40.38s", "+3.8%"]
    ]
    add_table(slide, table_data, Inches(0.8), Inches(2), Inches(8.4), Inches(4.5))
    
    # ============================================================================
    # SLIDE 7: Performance - Key Finding
    # ============================================================================
    slide = create_content_slide(prs, "Performance: Key Finding")
    
    # Big number
    big_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1.5))
    tf = big_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MicroVM Overhead: +3.8%"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER
    
    points = [
        "✓ Negligible performance impact for hardware-level isolation",
        "✓ Python hot starts actually FASTER on MicroVM (-2.7%)",
        "✓ Java shows minimal difference (-1.0% cold, +1.2% hot)",
        "✓ Language runtime initialization dominates (1.6-1.9× overhead)"
    ]
    add_bullet_points(slide, points, Inches(1), Inches(4.5), Inches(8), Inches(2))
    
    # ============================================================================
    # SLIDE 8: Cold vs Hot Start Analysis
    # ============================================================================
    slide = create_content_slide(prs, "Cold vs Warm Start Comparison")
    
    table_data = [
        ["Language", "Cluster", "Cold (s)", "Hot (s)", "Ratio"],
        ["Python", "C1", "41.84", "26.37", "1.6×"],
        ["Python", "C2", "47.17", "25.65", "1.8×"],
        ["JavaScript", "C1", "41.66", "24.72", "1.7×"],
        ["JavaScript", "C2", "45.43", "25.59", "1.8×"],
        ["Java", "C1", "45.17", "23.57", "1.9×"],
        ["Java", "C2", "44.73", "23.85", "1.9×"]
    ]
    add_table(slide, table_data, Inches(1.2), Inches(2), Inches(7.6), Inches(4.2))
    
    # ============================================================================
    # SLIDE 9: Energy Results - Overview
    # ============================================================================
    slide = create_content_slide(prs, "Energy Consumption: Overview")
    
    # Big numbers
    box1 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1.2))
    tf = box1.text_frame
    p = tf.paragraphs[0]
    p.text = "C1: 2812.4 J"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ENSEEIHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    box2 = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(1.2))
    tf = box2.text_frame
    p = tf.paragraphs[0]
    p.text = "C2: 2938.0 J"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Overhead
    box3 = slide.shapes.add_textbox(Inches(2.5), Inches(3.5), Inches(5), Inches(1))
    tf = box3.text_frame
    p = tf.paragraphs[0]
    p.text = "MicroVM Overhead: +125.6 J (+4.47%)"
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    points = [
        "Baseline idle consumption: 2800 J",
        "C1 overhead: +12.4 J (network I/O dominates)",
        "C2 overhead: +138.0 J (virtualization + network)"
    ]
    add_bullet_points(slide, points, Inches(1.5), Inches(5), Inches(7), Inches(1.5))
    
    # ============================================================================
    # SLIDE 10: Energy - THE KEY FINDING
    # ============================================================================
    slide = create_content_slide(prs, "Energy: The Surprising Discovery")
    
    # Highlight box
    big_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = big_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MicroVMs are 5× MORE EFFICIENT\nfor Warm Starts!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER
    
    table_data = [
        ["", "Cold → Warm Savings"],
        ["Cluster 1 (Container)", "0.5% (15.7 J)"],
        ["Cluster 2 (MicroVM)", "2.6% (77.3 J)"],
        ["MicroVM Advantage", "5× better efficiency"]
    ]
    add_table(slide, table_data, Inches(2), Inches(4.5), Inches(6), Inches(1.8))
    
    # ============================================================================
    # SLIDE 11: Energy by Language
    # ============================================================================
    slide = create_content_slide(prs, "Energy Consumption by Language")
    
    table_data = [
        ["Test", "C1 (J)", "C2 (J)", "Savings (%)"],
        ["Python Cold", "2812.8", "2972.6", "—"],
        ["Python Warm", "2800.4", "2897.8", "2.5%"],
        ["JavaScript Cold", "2814.2", "2973.0", "—"],
        ["JavaScript Warm", "2805.4", "2897.2", "2.5%"],
        ["Java Cold", "2833.6", "2984.2", "—"],
        ["Java Warm", "2807.8", "2903.0", "2.7%"]
    ]
    add_table(slide, table_data, Inches(1), Inches(2), Inches(8), Inches(4))
    
    # ============================================================================
    # SLIDE 12: Network Performance Summary
    # ============================================================================
    slide = create_content_slide(prs, "Network Performance")
    
    table_data = [
        ["Metric", "C1 (VM)", "C2 (µVM)", "Overhead"],
        ["TCP Latency", "7.07 ms", "8.09 ms", "+1.02 ms"],
        ["HTTP Latency", "266.30 ms", "255.51 ms", "-10.79 ms"],
        ["DNS Resolution", "7.77 ms", "8.26 ms", "+0.49 ms"],
        ["Bandwidth", "46.62 Mbps", "44.91 Mbps", "-1.71 Mbps"]
    ]
    add_table(slide, table_data, Inches(1.5), Inches(2.2), Inches(7), Inches(3))
    
    points = [
        "✓ TCP overhead: +1ms (negligible for 200-300ms HTTP requests)",
        "✓ HTTP performance: EQUIVALENT (variability dominates)",
        "✓ Bandwidth: -3.7% (< 300ms difference for 40MB transfer)"
    ]
    add_bullet_points(slide, points, Inches(1), Inches(5.5), Inches(8), Inches(1.5))
    
    # ============================================================================
    # SLIDE 13: Key Findings Summary
    # ============================================================================
    slide = create_content_slide(prs, "Key Findings")
    points = [
        "1. MicroVM Performance Overhead: Only +3.8%",
        "   → Minimal cost for hardware-level isolation",
        "",
        "2. Energy Paradox: MicroVMs 5× more efficient for warm starts",
        "   → Better for production workloads with container reuse",
        "",
        "3. Network Performance: Equivalent between platforms",
        "   → Virtualization doesn't compromise I/O capabilities",
        "",
        "4. Cold Start Dominates: Language choice > isolation method",
        "   → 1.6-1.9× overhead from runtime initialization"
    ]
    add_bullet_points(slide, points, Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    
    # ============================================================================
    # SLIDE 14: Production Implications
    # ============================================================================
    slide = create_content_slide(prs, "Production Implications")
    
    # Box 1: When to use MicroVMs
    box1 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(2.5))
    tf = box1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Use MicroVMs for:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    p = tf.add_paragraph()
    p.text = "• Multi-tenant platforms\n• Untrusted code\n• Compliance workloads\n• High-frequency functions\n• Security-critical applications"
    p.font.size = Pt(14)
    p.level = 0
    
    # Box 2: When to use Containers
    box2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(2.5))
    tf = box2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Use Containers for:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ENSEEIHT_BLUE
    
    p = tf.add_paragraph()
    p.text = "• Single-tenant environments\n• Trusted internal code\n• Resource-constrained nodes\n• Ultra-low latency (<1ms)\n• Development/testing"
    p.font.size = Pt(14)
    p.level = 0
    
    # Bottom recommendation
    rec_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    tf = rec_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Recommendation: Hybrid approach\nMicroVMs for security-sensitive + Containers for internal workloads"
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # ============================================================================
    # SLIDE 15: Optimization Recommendations
    # ============================================================================
    slide = create_content_slide(prs, "Optimization Strategies")
    points = [
        "Maximize Warm Starts (5× energy advantage for MicroVMs):",
        "  • Aggressive keep-alive policies (10+ min timeout)",
        "  • Provisioned concurrency for latency-sensitive functions",
        "  • Predictive pre-warming based on usage patterns",
        "",
        "Language Selection:",
        "  • Java: Best for long-running, frequent invocations",
        "  • JavaScript: Balanced performance, good general choice",
        "  • Python: Best for infrequent, short-lived functions",
        "",
        "At 99% warm starts: MicroVM overhead drops to only +3.4%"
    ]
    add_bullet_points(slide, points, Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    
    # ============================================================================
    # SLIDE 16: Contributions
    # ============================================================================
    slide = create_content_slide(prs, "Research Contributions")
    points = [
        "1. Comprehensive Benchmarking Framework",
        "   → Systematic evaluation across performance, energy, network",
        "",
        "2. First Energy Comparison of Container vs MicroVM",
        "   → Real-time RAPL monitoring with 30-second granularity",
        "",
        "3. Multi-Language Cross-Analysis",
        "   → Python, JavaScript, Java across isolation mechanisms",
        "",
        "4. Evidence-Based Recommendations",
        "   → Practical guidelines for FaaS architects",
        "",
        "5. Open-Source Implementation",
        "   → Reproducible benchmarks for community use"
    ]
    add_bullet_points(slide, points, Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    
    # ============================================================================
    # SLIDE 17: Limitations
    # ============================================================================
    slide = create_content_slide(prs, "Limitations & Future Work")
    
    # Limitations
    lim_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(2.5))
    tf = lim_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Limitations:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED
    
    p = tf.add_paragraph()
    p.text = "• Small sample size (n=5)\n• Single-node clusters\n• Network workloads only\n• Platform-specific (OpenWhisk)"
    p.font.size = Pt(14)
    
    # Future work - FIXED TYPO HERE
    future_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(2.5))
    tf = future_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Future Work:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    p = tf.add_paragraph()
    p.text = "• Larger scale (n≥30, multi-node)\n• CPU/memory workloads\n• Firecracker, gVisor comparison\n• ML-based optimization"
    p.font.size = Pt(14)
    
    # Impact
    impact_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(1.5))
    tf = impact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Despite limitations, findings provide strong evidence for\nMicroVM viability in production serverless platforms"
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # ============================================================================
    # SLIDE 18: Conclusion
    # ============================================================================
    slide = create_content_slide(prs, "Conclusion")
    
    # Main conclusion
    conclusion_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = conclusion_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MicroVMs offer hardware-level security\nat an acceptable efficiency cost"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ENSEEIHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    points = [
        "Trade-off is compelling:",
        "  • 4-5% efficiency investment",
        "  • Hardware-enforced isolation",
        "  • Regulatory compliance capability",
        "  • Superior warm start efficiency (5×)",
        "",
        "For multi-tenant FaaS platforms executing untrusted code,",
        "MicroVMs represent a production-ready, secure solution."
    ]
    add_bullet_points(slide, points, Inches(1), Inches(4), Inches(8), Inches(2.5))
    
    # ============================================================================
    # SLIDE 19: Thank You
    # ============================================================================
    slide = create_title_slide(
        prs,
        "Thank You",
        "Questions?"
    )
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ayoub SAMI - ayoub.sami@etu.inp-n7.fr\nENSEEIHT - Projet Long 2025-2026"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = "Network_Benchmark_Soutenance.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")

if __name__ == '__main__':
    main()